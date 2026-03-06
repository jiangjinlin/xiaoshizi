from pptx import Presentation
from typing import Dict, Any, List
import zipfile
import posixpath
import xml.etree.ElementTree as ET

NS = {
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
}

# 规范化 ZIP 内路径（OPC 部件名通常以“/”开头，但 zipfile 不包含前导“/”）
_def_lstrip = lambda s: s.lstrip('/') if isinstance(s, str) else s

def _norm_zip_path(p: str) -> str:
    return _def_lstrip(p)


def _shape_text(shape) -> str:
    if not hasattr(shape, "text_frame") or shape.text_frame is None:
        return ""
    return "\n".join(p.text or "" for p in shape.text_frame.paragraphs).strip()


# ============ OOXML 基础 ============

def _zip_open(file_path: str):
    return zipfile.ZipFile(file_path, 'r')


def _zip_read_xml(z: zipfile.ZipFile, path: str):
    try:
        norm = _norm_zip_path(path)
        with z.open(norm) as f:
            return ET.fromstring(f.read())
    except KeyError:
        return None


def _rels_path_for_part(partname: str) -> str:
    partname = _norm_zip_path(partname)
    dirname = posixpath.dirname(partname)
    base = posixpath.basename(partname)
    return posixpath.join(dirname, '_rels', base + '.rels')


def _find_rel_targets(z: zipfile.ZipFile, source_part: str, type_suffix: str) -> List[str]:
    rels_path = _rels_path_for_part(source_part)
    root = _zip_read_xml(z, rels_path)
    if root is None:
        return []
    targets = []
    for rel in root.findall('.//{http://schemas.openxmlformats.org/package/2006/relationships}Relationship'):
        rtype = rel.get('Type', '')
        if rtype.endswith(type_suffix):
            tgt = rel.get('Target')
            if not tgt:
                continue
            base_dir = posixpath.dirname(_norm_zip_path(source_part))
            abs_path = posixpath.normpath(posixpath.join(base_dir, tgt))
            targets.append(_norm_zip_path(abs_path))
    return targets


def _layout_for_slide(z: zipfile.ZipFile, slide_partname: str) -> str | None:
    targets = _find_rel_targets(z, slide_partname, '/slideLayout')
    return targets[0] if targets else None


def _master_for_layout(z: zipfile.ZipFile, layout_partname: str) -> str | None:
    targets = _find_rel_targets(z, layout_partname, '/slideMaster')
    return targets[0] if targets else None


def _theme_for_master(z: zipfile.ZipFile, master_partname: str) -> str | None:
    targets = _find_rel_targets(z, master_partname, '/theme')
    return targets[0] if targets else None


def _layout_type(z: zipfile.ZipFile, layout_partname: str) -> str | None:
    root = _zip_read_xml(z, layout_partname)
    if root is None:
        return None
    t = root.get('type')
    if t:
        return t
    node = root.find('p:sldLayout', NS)
    return node.get('type') if node is not None else None


def _layout_has_two_bodies(z: zipfile.ZipFile, layout_partname: str) -> bool:
    root = _zip_read_xml(z, layout_partname)
    if root is None:
        return False
    bodies = root.findall('.//p:ph[@type="body"]', NS)
    return len(bodies) >= 2


def _theme_name_contains(z: zipfile.ZipFile, keyword_list: List[str]) -> bool:
    for name in z.namelist():
        if name.startswith('ppt/theme/') and name.endswith('.xml'):
            root = _zip_read_xml(z, name)
            if root is None:
                continue
            tname = root.get('name') or ''
            if any(k in tname for k in keyword_list):
                return True
    return False


# ============ 文本与形状/图片位置信息 ============

def _build_parent_map(root: ET.Element) -> Dict[ET.Element, ET.Element]:
    pm: Dict[ET.Element, ET.Element] = {}
    for parent in root.iter():
        for child in list(parent):
            pm[child] = parent
    return pm


def _grp_offset_x(el: ET.Element, pm: Dict[ET.Element, ET.Element]) -> int:
    total = 0
    p = pm.get(el)
    while p is not None:
        # 分组 p:grpSp 的偏移在 p:grpSpPr/a:xfrm/a:off
        if p.tag.endswith('grpSp'):
            gxfrm = p.find('p:grpSpPr/a:xfrm', NS)
            if gxfrm is not None:
                off = gxfrm.find('a:off', NS)
                if off is not None and off.get('x') is not None:
                    try:
                        total += int(off.get('x'))
                    except Exception:
                        pass
        p = pm.get(p)
    return total


def _collect_shapes_info(root: ET.Element) -> List[Dict[str, Any]]:
    pm = _build_parent_map(root)
    infos: List[Dict[str, Any]] = []
    # 普通形状 p:sp（位置在 p:spPr/a:xfrm）
    for sp in root.findall('.//p:sp', NS):
        ph = sp.find('.//p:nvSpPr/p:nvPr/p:ph', NS)
        ph_type = ph.get('type') if ph is not None else None
        ph_idx = None
        if ph is not None and ph.get('idx') is not None:
            try:
                ph_idx = int(ph.get('idx'))
            except Exception:
                ph_idx = None
        xfrm = sp.find('p:spPr/a:xfrm', NS)
        off = xfrm.find('a:off', NS) if xfrm is not None else None
        ext = xfrm.find('a:ext', NS) if xfrm is not None else None
        x = int(off.get('x')) if off is not None and off.get('x') is not None else 0
        w = int(ext.get('cx')) if ext is not None and ext.get('cx') is not None else 0
        x += _grp_offset_x(sp, pm)
        texts = [t.text for t in sp.findall('.//a:t', NS) if t.text]
        infos.append({'tag': 'sp', 'ph_type': ph_type, 'ph_idx': ph_idx, 'x': x, 'w': w, 'texts': texts, 'el': sp})
    # 图片 p:pic（位置在 p:spPr/a:xfrm）
    for pic in root.findall('.//p:pic', NS):
        xfrm = pic.find('p:spPr/a:xfrm', NS)
        off = xfrm.find('a:off', NS) if xfrm is not None else None
        ext = xfrm.find('a:ext', NS) if xfrm is not None else None
        x = int(off.get('x')) if off is not None and off.get('x') is not None else 0
        w = int(ext.get('cx')) if ext is not None and ext.get('cx') is not None else 0
        x += _grp_offset_x(pic, pm)
        cNvPr = pic.find('.//p:nvPicPr/p:cNvPr', NS)
        name = cNvPr.get('name') if cNvPr is not None else ''
        descr = cNvPr.get('descr') if cNvPr is not None else ''
        infos.append({'tag': 'pic', 'x': x, 'w': w, 'name': name, 'descr': descr, 'el': pic})
    # 图形框 p:graphicFrame（位置在 a:xfrm）
    for gf in root.findall('.//p:graphicFrame', NS):
        xfrm = gf.find('a:xfrm', NS)
        off = xfrm.find('a:off', NS) if xfrm is not None else None
        ext = xfrm.find('a:ext', NS) if xfrm is not None else None
        x = int(off.get('x')) if off is not None and off.get('x') is not None else 0
        w = int(ext.get('cx')) if ext is not None and ext.get('cx') is not None else 0
        x += _grp_offset_x(gf, pm)
        texts = [t.text for t in gf.findall('.//a:t', NS) if t.text]
        infos.append({'tag': 'gfr', 'x': x, 'w': w, 'texts': texts, 'el': gf})
    return infos


def _right_body_center_x(slide_root: ET.Element) -> int | None:
    infos = _collect_shapes_info(slide_root)
    bodies = [i for i in infos if i.get('tag') == 'sp' and i.get('ph_type') == 'body']
    if len(bodies) < 1:
        return None
    right = max(bodies, key=lambda i: i.get('x', 0))
    return right['x'] + right.get('w', 0) // 2


def _slide2_has_pic_in_right(slide_root: ET.Element) -> bool:
    cx = _right_body_center_x(slide_root)
    if cx is None:
        return False
    for info in _collect_shapes_info(slide_root):
        if info.get('tag') == 'pic':
            pc = info.get('x', 0) + info.get('w', 0) // 2
            if pc >= cx:
                return True
    return False


# ============ 页脚、页码、切换、放映、动画、样式 ============

def _slide_has_footer_text(z: zipfile.ZipFile, slide_partname: str, text: str) -> bool:
    root = _zip_read_xml(z, slide_partname)
    if root is None:
        return False
    # 先查页脚占位符；若未找到则查整页文本
    for sp in root.findall('.//p:sp', NS):
        ph = sp.find('.//p:ph', NS)
        if ph is not None and ph.get('type') == 'ftr':
            t_list = [t.text or '' for t in sp.findall('.//a:t', NS)]
            if text in ''.join(t_list):
                return True
    all_text = ''.join([t.text or '' for t in root.findall('.//a:t', NS)])
    return text in all_text


def _layout_or_master_has_footer_text(z: zipfile.ZipFile, layout_partname: str | None, text: str) -> bool:
    if not layout_partname:
        return False
    # 版式
    lroot = _zip_read_xml(z, layout_partname)
    if lroot is not None:
        for sp in lroot.findall('.//p:sp', NS):
            ph = sp.find('.//p:ph', NS)
            if ph is not None and ph.get('type') == 'ftr':
                t_list = [t.text or '' for t in sp.findall('.//a:t', NS)]
                if text in ''.join(t_list):
                    return True
    # 母版
    master = _master_for_layout(z, layout_partname)
    if master:
        mroot = _zip_read_xml(z, master)
        if mroot is not None:
            for sp in mroot.findall('.//p:sp', NS):
                ph = sp.find('.//p:ph', NS)
                if ph is not None and ph.get('type') == 'ftr':
                    t_list = [t.text or '' for t in sp.findall('.//a:t', NS)]
                    if text in ''.join(t_list):
                        return True
    return False


def _slide_or_anc_has_sldnum(z: zipfile.ZipFile, slide_partname: str) -> bool:
    root = _zip_read_xml(z, slide_partname)
    # 方案A：占位符 sldNum（常规）
    if root is not None and root.find('.//p:ph[@type="sldNum"]', NS) is not None:
        return True
    # 方案B：文本字段 a:fld type="slidenum"（Office 2010 常见）
    if root is not None and root.find('.//a:fld[@type="slidenum"]', NS) is not None:
        return True
    layout = _layout_for_slide(z, slide_partname)
    if layout:
        lroot = _zip_read_xml(z, layout)
        if lroot is not None and (lroot.find('.//p:ph[@type="sldNum"]', NS) is not None or lroot.find('.//a:fld[@type="slidenum"]', NS) is not None):
            return True
        master = _master_for_layout(z, layout)
        if master:
            mroot = _zip_read_xml(z, master)
            if mroot is not None and (mroot.find('.//p:ph[@type="sldNum"]', NS) is not None or mroot.find('.//a:fld[@type="slidenum"]', NS) is not None):
                return True
    return False


def _all_slides_transition_wipe(z: zipfile.ZipFile, slide_partnames: List[str]) -> bool:
    for spath in slide_partnames:
        root = _zip_read_xml(z, spath)
        if root is None:
            return False
        trans = root.find('.//p:transition', NS)
        if trans is None or trans.find('p:wipe', NS) is None:
            return False
    return True


def _show_mode_browse(z: zipfile.ZipFile) -> bool:
    pres = _zip_read_xml(z, 'ppt/presentation.xml')
    if pres is not None:
        sp = pres.find('p:showPr', NS)
        if sp is not None and ((sp.find('p:browse', NS) is not None) or (sp.find('p:kiosk', NS) is not None)):
            return True
    # 兼容：有些版本在 presProps.xml 中保存显示属性
    pres_props = _zip_read_xml(z, 'ppt/presProps.xml')
    if pres_props is not None:
        sp = pres_props.find('p:showPr', NS)
        if sp is not None and ((sp.find('p:browse', NS) is not None) or (sp.find('p:kiosk', NS) is not None)):
            return True
    return False


def _slide2_pic_has_soft_edge(z: zipfile.ZipFile, slide_partname: str) -> bool:
    root = _zip_read_xml(z, slide_partname)
    if root is None:
        return False
    for pic in root.findall('.//p:pic', NS):
        eff = pic.find('.//a:effectLst', NS)
        if eff is not None and eff.find('a:softEdge', NS) is not None:
            return True
    return False


def _slide_has_floatin_anim_from_bottom(z: zipfile.ZipFile, slide_partname: str) -> bool:
    root = _zip_read_xml(z, slide_partname)
    if root is None:
        return False
    timing = root.find('.//p:timing', NS)
    if timing is None:
        return False
    # 1) 进入类动画：animEffect transition="in"（filter 任意，如 fade/fly/float）
    has_in = False
    for ae in timing.findall('.//p:animEffect', NS):
        trans = (ae.get('transition') or '').lower()
        if 'in' in trans:
            has_in = True
            break
    # 2) 方向：检测纵向位移动画 ppt_y，从偏移到基准（如 #ppt_y-.1 -> #ppt_y 或 #ppt_y+.1 -> #ppt_y）
    def _has_vertical_move_down_or_up() -> bool:
        for anim in timing.findall('.//p:anim', NS):
            # 仅关注 y 方向
            attr_ok = False
            for an in anim.findall('.//p:attrName', NS):
                if (an.text or '').strip() == 'ppt_y':
                    attr_ok = True
                    break
            if not attr_ok:
                continue
            # 读取 tavLst 的首尾值
            vals: List[str] = []
            for tav in anim.findall('.//p:tav', NS):
                v = tav.find('.//p:strVal', NS)
                if v is not None and v.get('val'):
                    vals.append(v.get('val'))
            if len(vals) >= 2:
                first = vals[0].lower()
                last = vals[-1].lower()
                if first != last and '#ppt_y' in first and '#ppt_y' in last:
                    # 认为存在从偏移恢复到基准的纵向位移
                    return True
        return False
    has_vert = _has_vertical_move_down_or_up()

    # 3) 预设进入动画方向（保留原有识别）
    if not (has_in and has_vert):
        for ae in timing.findall('.//p:animEffect', NS):
            filt = (ae.get('filter') or '').lower()
            trans = (ae.get('transition') or '').lower()
            dirv = (ae.get('dir') or '').lower()
            if (('fly' in filt or 'float' in filt) and ('in' in trans or 'in' in filt)) and (dirv in ('b', 'bottom', 'down')):
                return True
        for elem in timing.iter():
            pc = (elem.get('presetClass') or '').lower()
            ps = (elem.get('presetSubtype') or '').lower()
            if pc == 'entr' and ps in ('down','bottom','frombottom','bottom-left','bottom-right'):
                return True
    # 4) 宽松兜底：属性文本匹配
    if has_in and has_vert:
        return True
    for elem in timing.iter():
        attrs = {k.split('}')[-1]: v for k, v in elem.attrib.items()}
        text = ' '.join(list(attrs.values())).lower()
        if ('float' in text or 'fly' in text or 'in' in text) and ('down' in text or 'bottom' in text or attrs.get('dir','').lower() in ('b','bottom','down')):
            return True
    return False


def _get_pres_center_x(z: zipfile.ZipFile) -> int | None:
    pres = _zip_read_xml(z, 'ppt/presentation.xml')
    if pres is None:
        return None
    sldSz = pres.find('p:sldSz', NS)
    if sldSz is None:
        return None
    try:
        cx = int(sldSz.get('cx'))
        return cx // 2
    except Exception:
        return None


def _infer_two_columns_by_positions(slide_root: ET.Element, center_x: int | None = None) -> bool:
    infos = _collect_shapes_info(slide_root)
    # 任何包含文本的形状或图片/图形框都视为有效内容
    def is_content(i: Dict[str, Any]) -> bool:
        if i.get('tag') in ('pic', 'gfr'):
            return True
        return bool(''.join(i.get('texts') or []).strip())
    content_shapes = [i for i in infos if i.get('tag') in ('sp','pic','gfr') and is_content(i)]
    if not content_shapes:
        return False
    if center_x is None:
        minx = min(i.get('x', 0) for i in content_shapes)
        maxx = max(i.get('x', 0) + i.get('w', 0) for i in content_shapes)
        center_x = (minx + maxx) // 2
    left = any(i.get('x', 0) + i.get('w', 0) // 2 < center_x for i in content_shapes)
    right = any(i.get('x', 0) + i.get('w', 0) // 2 >= center_x for i in content_shapes)
    return left and right


def _any_pic_right_of(slide_root: ET.Element, center_x: int) -> bool:
    for info in _collect_shapes_info(slide_root):
        if info.get('tag') == 'pic':
            pc = info.get('x', 0) + info.get('w', 0) // 2
            if pc >= center_x:
                return True
    return False


def _content_shapes(slide_root: ET.Element) -> List[Dict[str, Any]]:
    infos = _collect_shapes_info(slide_root)
    contents: List[Dict[str, Any]] = []
    for i in infos:
        tag = i.get('tag')
        if tag in ('pic', 'gfr'):
            contents.append(i)
        elif tag == 'sp':
            ph_type = i.get('ph_type')
            if ph_type in ('title', 'ftr', 'sldNum'):
                continue
            if ''.join(i.get('texts') or []).strip():
                contents.append(i)
    return contents


# ============ 评分主流程 ============

def score_ppt(file_path: str) -> tuple[int, int, str]:
    prs = Presentation(file_path)
    slides = prs.slides
    slide_count = len(slides)

    details: List[Dict[str, Any]] = []

    # 分值配置（总计10分）
    # 1 主题 1.0
    # 2 第1张版式/标题/字体/字号 2.0
    # 3 第2张两栏/标题/右栏图片 2.0
    # 4 图片样式+进入动画 1.0
    # 5 页码+页脚 1.5
    # 6 第3张两栏+右栏有内容+标题+删除第4张 1.5
    # 7 切换+放映方式 1.0

    total = 10.0
    score = 0.0

    # 打开包与部件路径
    z = _zip_open(file_path)
    try:
        slide_partnames = [s.part.partname.replace('\\', '/') for s in slides]
    except Exception:
        slide_partnames = [f'/ppt/slides/slide{i+1}.xml' for i in range(slide_count)]

    # 读取各页XML
    slide_roots: List[ET.Element] = []
    for p in slide_partnames:
        slide_roots.append(_zip_read_xml(z, p))

    # 1) 主题“奥斯汀” 1.0
    theme_ok = _theme_name_contains(z, ['奥斯汀', 'Austin'])
    details.append({'item': '主题为“奥斯汀”', 'ok': bool(theme_ok), 'score': 1.0 if theme_ok else 0.0, 'max': 1.0, 'note': ''})
    if theme_ok:
        score += 1.0

    # 2) 第1张 2.0
    if slide_count >= 1 and slide_roots[0] is not None:
        s1_part = slide_partnames[0]
        s1_layout = _layout_for_slide(z, s1_part)
        s1_layout_type = _layout_type(z, s1_layout) if s1_layout else None
        is_title_layout = (s1_layout_type == 'title')
        details.append({'item': '第1张：版式为标题幻灯片', 'ok': bool(is_title_layout), 'score': 0.5 if is_title_layout else 0.0, 'max': 0.5, 'note': f'layout.type={s1_layout_type or "?"}'})
        if is_title_layout:
            score += 0.5

        s1 = slide_roots[0]
        # 提取标题/副标题文本与字体/字号（取第一个匹配占位符的第一段第���run为准）
        def _extract_text_font_size(sp: ET.Element):
            txt = ''.join([t.text for t in sp.findall('.//a:t', NS) if t.text])
            rPr = sp.find('.//a:r/a:rPr', NS)
            sz = rPr.get('sz') if rPr is not None else None  # 1/100 pt
            face = ''
            if rPr is not None:
                ea = rPr.find('a:ea', NS)
                lat = rPr.find('a:latin', NS)
                if ea is not None and ea.get('typeface'):
                    face = ea.get('typeface')
                elif lat is not None and lat.get('typeface'):
                    face = lat.get('typeface')
            return txt, sz, face

        # 找占位符
        title_sp = None
        sub_sp = None
        for sp in s1.findall('.//p:sp', NS):
            ph = sp.find('.//p:ph', NS)
            if ph is None:
                continue
            if ph.get('type') == 'title' and title_sp is None:
                title_sp = sp
            if ph.get('type') == 'subTitle' and sub_sp is None:
                sub_sp = sp
        # 回退
        if title_sp is None:
            title_sp = s1.find('.//p:sp', NS)
        title_text, title_sz, title_face = _extract_text_font_size(title_sp) if title_sp is not None else ('', None, '')
        sub_text, sub_sz, sub_face = _extract_text_font_size(sub_sp) if sub_sp is not None else ('', None, '')

        # 标题文本
        title_text_ok = (title_text.strip() == '草原生态补奖政策实施十年')
        details.append({'item': '第1张：标题文本为“草原生态补奖政策实施十年”', 'ok': bool(title_text_ok), 'score': 0.5 if title_text_ok else 0.0, 'max': 0.5, 'note': f"实际='{title_text}'"})
        if title_text_ok:
            score += 0.5
        # 标题字体/字号
        title_face_ok = (title_face in ('华文行楷', 'STXingkai'))
        details.append({'item': '第1张：标题字体为“华文行楷”', 'ok': bool(title_face_ok), 'score': 0.5 if title_face_ok else 0.0, 'max': 0.5, 'note': f"检测='{title_face or ''}'"})
        if title_face_ok:
            score += 0.5
        title_size_ok = (title_sz is not None and title_sz.isdigit() and abs(int(title_sz) - 4000) <= 50)
        details.append({'item': '第1张：标题字号为40号', 'ok': bool(title_size_ok), 'score': 0.25 if title_size_ok else 0.0, 'max': 0.25, 'note': f"检测='{title_sz or ''}'(1/100pt)"})
        if title_size_ok:
            score += 0.25
        # 副标题字体/字号
        sub_face_ok = (sub_face in ('华文楷体', 'STKaiti'))
        details.append({'item': '第1张：副标题字体为“华文楷体”', 'ok': bool(sub_face_ok), 'score': 0.15 if sub_face_ok else 0.0, 'max': 0.15, 'note': f"检测='{sub_face or ''}'"})
        if sub_face_ok:
            score += 0.15
        sub_size_ok = (sub_sz is not None and sub_sz.isdigit() and abs(int(sub_sz) - 2000) <= 50)
        details.append({'item': '第1张：副标题字号为20号', 'ok': bool(sub_size_ok), 'score': 0.10 if sub_size_ok else 0.0, 'max': 0.10, 'note': f"检测='{sub_sz or ''}'(1/100pt)"})
        if sub_size_ok:
            score += 0.10

    # 3) 第2张 两栏/标题/右栏图片 2.0
    if slide_count >= 2 and slide_roots[1] is not None:
        s2_part = slide_partnames[1]
        s2_layout = _layout_for_slide(z, s2_part)
        two_col_ok = _layout_has_two_bodies(z, s2_layout) if s2_layout else False
        if not two_col_ok:
            # 退化到实际页内占位符
            infos = _collect_shapes_info(slide_roots[1])
            bodies = [i for i in infos if i.get('tag') == 'sp' and i.get('ph_type') == 'body']
            two_col_ok = len(bodies) >= 2
        if not two_col_ok:
            # 再次兜底：基于文本/图片/图形框分布推断两栏
            cx_pres = _get_pres_center_x(z)
            two_col_ok = _infer_two_columns_by_positions(slide_roots[1], cx_pres)
        if not two_col_ok:
            # 最终兜底：只要内容形状数量≥2（排除标题/页脚/页码），视为两栏
            two_col_ok = len(_content_shapes(slide_roots[1])) >= 2
        details.append({'item': '第2张：两栏内容版式', 'ok': bool(two_col_ok), 'score': 0.8 if two_col_ok else 0.0, 'max': 0.8, 'note': ''})
        if two_col_ok:
            score += 0.8
        # 标题
        title_ok = False
        s2 = slide_roots[1]
        for sp in s2.findall('.//p:sp', NS):
            ph = sp.find('.//p:ph', NS)
            if ph is not None and ph.get('type') in ('title',) or (ph is not None and ph.get('idx') == '0'):
                s2_title = ''.join([t.text for t in sp.findall('.//a:t', NS) if t.text])
                title_ok = (s2_title.strip() == '万马奔腾')
                break
        details.append({'item': '第2张：标题为“万马奔腾”', 'ok': bool(title_ok), 'score': 0.4 if title_ok else 0.0, 'max': 0.4, 'note': ''})
        if title_ok:
            score += 0.4
        # 图片与位置
        pics = slide_roots[1].findall('.//p:pic', NS)
        pic_ok = len(pics) > 0
        details.append({'item': '第2张：插入图片', 'ok': bool(pic_ok), 'score': 0.5 if pic_ok else 0.0, 'max': 0.5, 'note': ''})
        if pic_ok:
            score += 0.5
        in_right = False
        if pic_ok:
            # 优先依据右侧正文中心；若不可得，退化为版面中心
            right_cx = _right_body_center_x(slide_roots[1])
            if right_cx is None:
                right_cx = _get_pres_center_x(z)
            if right_cx is None:
                # 最终兜底：用当前页形状范围居中（含图片/图形框）
                infos = _collect_shapes_info(slide_roots[1])
                xs = [i.get('x', 0) for i in infos if i.get('tag') in ('sp','pic','gfr')]
                xe = [i.get('x', 0) + i.get('w', 0) for i in infos if i.get('tag') in ('sp','pic','gfr')]
                if xs and xe:
                    right_cx = (min(xs) + max(xe)) // 2
            if right_cx is not None:
                in_right = _any_pic_right_of(slide_roots[1], right_cx)
        details.append({'item': '第2张：图片位于右栏', 'ok': bool(in_right), 'score': 0.3 if in_right else 0.0, 'max': 0.3, 'note': ''})
        if in_right:
            score += 0.3

    # 4) 图片样式与动画 1.0
    if slide_count >= 2 and slide_roots[1] is not None:
        soft_ok = _slide2_pic_has_soft_edge(z, slide_partnames[1])
        details.append({'item': '第2张：图片样式“柔化边缘”', 'ok': bool(soft_ok), 'score': 0.5 if soft_ok else 0.0, 'max': 0.5, 'note': ''})
        if soft_ok:
            score += 0.5
        anim_ok = _slide_has_floatin_anim_from_bottom(z, slide_partnames[1])
        details.append({'item': '第2张：图片进入动画“浮入-下方”', 'ok': bool(anim_ok), 'score': 0.5 if anim_ok else 0.0, 'max': 0.5, 'note': ''})
        if anim_ok:
            score += 0.5

    # 5) 页码与页脚 1.5
    footer_text = '我的草原我的马'
    footer_all = True
    sldnum_all = True
    for i, part in enumerate(slide_partnames):
        lay = _layout_for_slide(z, part)
        has_footer = _slide_has_footer_text(z, part, footer_text) or _layout_or_master_has_footer_text(z, lay, footer_text)
        if not has_footer:
            footer_all = False
        if not _slide_or_anc_has_sldnum(z, part):
            sldnum_all = False
    details.append({'item': '所有幻灯片：包含页脚“我的草原我的马”', 'ok': bool(footer_all), 'score': 0.8 if footer_all else 0.0, 'max': 0.8, 'note': ''})
    if footer_all:
        score += 0.8
    details.append({'item': '所有幻灯片：包含页码', 'ok': bool(sldnum_all), 'score': 0.7 if sldnum_all else 0.0, 'max': 0.7, 'note': ''})
    if sldnum_all:
        score += 0.7

    # 6) 第3张：两栏+右栏有内容+标题+删除第4张 1.5
    if slide_count >= 3 and slide_roots[2] is not None:
        s3_part = slide_partnames[2]
        s3_layout = _layout_for_slide(z, s3_part)
        s3_two = _layout_has_two_bodies(z, s3_layout) if s3_layout else False
        if not s3_two:
            infos3 = _collect_shapes_info(slide_roots[2])
            s3_two = len([i for i in infos3 if i.get('tag') == 'sp' and i.get('ph_type') == 'body']) >= 2
        if not s3_two:
            # 再次兜底：基于文本/图片/图形框分布推断两栏
            cx_pres3 = _get_pres_center_x(z)
            s3_two = _infer_two_columns_by_positions(slide_roots[2], cx_pres3)
        if not s3_two:
            # 最终兜底：内容形状数量≥2 即视作两栏
            s3_two = len(_content_shapes(slide_roots[2])) >= 2
        details.append({'item': '第3张：两栏内容版式', 'ok': bool(s3_two), 'score': 0.4 if s3_two else 0.0, 'max': 0.4, 'note': ''})
        if s3_two:
            score += 0.4
        # 标题
        t3_ok = False
        s3 = slide_roots[2]
        for sp in s3.findall('.//p:sp', NS):
            ph = sp.find('.//p:ph', NS)
            if ph is not None and ph.get('type') in ('title',) or (ph is not None and ph.get('idx') == '0'):
                s3_title = ''.join([t.text for t in sp.findall('.//a:t', NS) if t.text])
                t3_ok = (s3_title.strip() == '草原生态补奖政策')
                break
        details.append({'item': '第3张：标题为“草原生态补奖政策”', 'ok': bool(t3_ok), 'score': 0.4 if t3_ok else 0.0, 'max': 0.4, 'note': ''})
        if t3_ok:
            score += 0.4
        # 右栏有内容：右侧正文占位符含文本或右半区域存在图片
        s3_right_cx = _right_body_center_x(slide_roots[2])
        if s3_right_cx is None:
            s3_right_cx = _get_pres_center_x(z)
            if s3_right_cx is None:
                infos = _collect_shapes_info(slide_roots[2])
                xs = [i.get('x', 0) for i in infos if i.get('tag') in ('sp','pic','gfr')]
                xe = [i.get('x', 0) + i.get('w', 0) for i in infos if i.get('tag') in ('sp','pic','gfr')]
                if xs and xe:
                    s3_right_cx = (min(xs) + max(xe)) // 2
        right_has_content = False
        if s3_right_cx is not None:
            # 文本
            for info in _collect_shapes_info(slide_roots[2]):
                if info.get('tag') == 'sp' and (info.get('ph_type') == 'body' or ''.join(info.get('texts') or []).strip()):
                    cx = info.get('x', 0) + info.get('w', 0) // 2
                    if cx >= s3_right_cx and ''.join(info.get('texts') or []).strip():
                        right_has_content = True
                        break
            # 图形框或图片兜底
            if not right_has_content:
                infos2 = _collect_shapes_info(slide_roots[2])
                for info in infos2:
                    if info.get('tag') in ('pic','gfr'):
                        cx = info.get('x', 0) + info.get('w', 0) // 2
                        if cx >= s3_right_cx:
                            right_has_content = True
                            break
            # 坐标缺失兜底：所有内容形状坐标为0时，按“至少两块内容”或占位符idx顺序判断右侧
            if not right_has_content:
                contents = _content_shapes(slide_roots[2])
                if contents:
                    all_zero_pos = all((i.get('x',0)==0 and i.get('w',0)==0) for i in contents)
                    if all_zero_pos:
                        # 优先用占位符idx：最大idx视为右栏且需有文本；否则按数量≥2
                        spcs = [i for i in contents if i.get('tag')=='sp']
                        idxs = [i.get('ph_idx') for i in spcs if isinstance(i.get('ph_idx'), int)]
                        if len(set(idxs)) >= 2:
                            max_idx = max(idxs)
                            right_has_content = any(i.get('tag')=='sp' and i.get('ph_idx')==max_idx and ''.join(i.get('texts') or []).strip() for i in spcs)
                        else:
                            right_has_content = len(contents) >= 2
        else:
            # 无法得到右侧��心：以内容形状数量兜底
            right_has_content = len(_content_shapes(slide_roots[2])) >= 2
        details.append({'item': '第3张：右栏含内容', 'ok': bool(right_has_content), 'score': 0.4 if right_has_content else 0.0, 'max': 0.4, 'note': ''})
        if right_has_content:
            score += 0.4
    # 删除第4张
    no_slide4 = (slide_count == 3)
    details.append({'item': '删除第4张幻灯片（总页数为3）', 'ok': bool(no_slide4), 'score': 0.3 if no_slide4 else 0.0, 'max': 0.3, 'note': f'当前页数={slide_count}'})
    if no_slide4:
        score += 0.3

    # 7) 切换+放映方式 1.0
    wipe_all = _all_slides_transition_wipe(z, slide_partnames) if slide_partnames else False
    details.append({'item': '所有幻灯片：切换为“擦除”', 'ok': bool(wipe_all), 'score': 0.5 if wipe_all else 0.0, 'max': 0.5, 'note': ''})
    if wipe_all:
        score += 0.5
    browse_ok = _show_mode_browse(z)
    details.append({'item': '放映方式：观众自行浏览', 'ok': bool(browse_ok), 'score': 0.5 if browse_ok else 0.0, 'max': 0.5, 'note': ''})
    if browse_ok:
        score += 0.5

    # 关闭包
    z.close()

    # 组装反馈信息，参考 word.py 风格（每个小点5分）
    lines: List[str] = []
    per_point = 5
    items_cnt = 0
    pass_cnt = 0
    for d in details:
        item = d.get('item', '')
        ok = bool(d.get('ok', False))
        note = d.get('note', '')
        # 统一按小点计分
        items_cnt += 1
        gained = per_point if ok else 0
        if ok:
            pass_cnt += 1
        lines.append(f"通过：{item}（+{gained:g}分）" if ok else f"未通过：{item}（0分）")
        if note:
            lines.append(f"说明：{note}")

    # 计算与返回三元组（每个小点5分）
    total = items_cnt * per_point
    final_score = pass_cnt * per_point
    msg = '\n'.join(lines) if lines else '评分完成'
    return final_score, total, msg
